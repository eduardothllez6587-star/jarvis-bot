import os
import re
import json
import logging
import requests
import tempfile
from pathlib import Path
from datetime import datetime

import telebot

logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")
logger = logging.getLogger(__name__)

TELEGRAM_TOKEN    = os.environ["TELEGRAM_TOKEN"]
GROQ_API_KEY      = os.environ["GROQ_API_KEY"]
ELEVENLABS_API_KEY = os.environ["ELEVENLABS_API_KEY"]
ELEVENLABS_VOICE  = os.environ.get("ELEVENLABS_VOICE_ID", "4tRn1lSkEn13EVTuqb0g")
MEMORY_FILE       = "memory.json"
MODEL             = "llama-3.1-8b-instant"
GROQ_URL          = "https://api.groq.com/openai/v1/chat/completions"
GROQ_AUDIO_URL    = "https://api.groq.com/openai/v1/audio/transcriptions"

bot = telebot.TeleBot(TELEGRAM_TOKEN)

# ── Memoria ───────────────────────────────────────────────────────────────────
def load_memory():
    if Path(MEMORY_FILE).exists():
        try:
            return json.loads(Path(MEMORY_FILE).read_text())
        except Exception:
            pass
    return {}

def save_memory(data):
    Path(MEMORY_FILE).write_text(json.dumps(data, ensure_ascii=False, indent=2))

def get_user_memory(user_id):
    mem = load_memory()
    return mem.get(str(user_id), {"items": [], "history": []})

def update_user_memory(user_id, user_mem):
    mem = load_memory()
    mem[str(user_id)] = user_mem
    save_memory(mem)

def extract_memory_tag(text):
    match = re.search(r"<<<MEMORIA:(\{.*?\})>>>", text, re.DOTALL)
    if match:
        try:
            item = json.loads(match.group(1))
            clean = text[:match.start()].strip() + text[match.end():].strip()
            return clean.strip(), item
        except Exception:
            pass
    return text, None

# ── Detección de agente ───────────────────────────────────────────────────────
def detect_agent(text):
    t = text.lower()
    if any(w in t for w in ["imagen", "foto", "dibuja", "genera", "pinta", "anime", "ilustra"]):
        return "imagen"
    if any(w in t for w in ["codigo", "código", "python", "programar", "script", "funcion", "función", "bug"]):
        return "codigo"
    if any(w in t for w in ["escribe", "historia", "cuento", "poema", "redacta", "carta", "correo", "email"]):
        return "escribe"
    if any(w in t for w in ["word", "documento", ".docx"]):
        return "word"
    if any(w in t for w in ["excel", "hoja", "tabla", ".xlsx"]):
        return "excel"
    if any(w in t for w in ["powerpoint", "presentacion", "presentación", "slides", ".pptx"]):
        return "pptx"
    return "director"

# ── Prompts ───────────────────────────────────────────────────────────────────
SYSTEM_PROMPTS = {
    "director": (
        "Eres Jarvis, el asistente personal elite de Eduardo. "
        "Estilo J.A.R.V.I.S. de Tony Stark: inteligente, directo, ocasionalmente ingenioso, nunca verboso. "
        "Siempre responde en español. Respuestas cortas y directas para mensajes de voz. "
        "Si el usuario menciona proyectos, tareas o preferencias importantes, termina tu respuesta con: "
        '<<<MEMORIA:{"tipo":"proyecto","dato":"descripcion breve"}>>>'
    ),
    "codigo": "Eres Jarvis en modo ingeniero senior. Codigo limpio, explica brevemente. Responde en español.",
    "escribe": "Eres Jarvis en modo escritor maestro. Cualquier genero, con alma. Responde en español.",
    "imagen": (
        "Eres Jarvis en modo director visual. El usuario quiere una imagen. "
        "Devuelve SOLO un prompt optimizado en ingles para generacion de imagen. Sin explicaciones."
    ),
    "word": (
        "Eres Jarvis en modo asistente de documentos. Genera contenido completo en español. "
        "Usa # para titulos principales y ## para subtitulos."
    ),
    "excel": (
        "Eres Jarvis en modo analista. Devuelve los datos en formato CSV limpio con encabezados. "
        "Responde SOLO con el CSV."
    ),
    "pptx": (
        "Eres Jarvis en modo director de presentaciones. "
        'Genera el contenido en JSON: [{"titulo":"Slide 1","puntos":["punto 1","punto 2"]}] '
        "Devuelve SOLO el JSON."
    ),
}

# ── Groq texto ────────────────────────────────────────────────────────────────
def ask_groq(agent, user_message, history, memory_items):
    system = SYSTEM_PROMPTS.get(agent, SYSTEM_PROMPTS["director"])
    if memory_items:
        system += "\n\nMemoria previa de Eduardo:\n"
        system += "\n".join(f"- [{m.get('tipo','')}] {m.get('dato','')}" for m in memory_items)
    messages = [{"role": "system", "content": system}]
    messages += history[-10:]
    messages.append({"role": "user", "content": user_message})
    response = requests.post(
        GROQ_URL,
        headers={"Authorization": f"Bearer {GROQ_API_KEY}", "Content-Type": "application/json"},
        json={"model": MODEL, "messages": messages, "temperature": 0.7, "max_tokens": 2048},
        timeout=30
    )
    data = response.json()
    return data["choices"][0]["message"]["content"].strip()

# ── Groq audio (transcripción) ────────────────────────────────────────────────
def transcribe_audio(audio_path):
    with open(audio_path, "rb") as f:
        response = requests.post(
            GROQ_AUDIO_URL,
            headers={"Authorization": f"Bearer {GROQ_API_KEY}"},
            files={"file": (Path(audio_path).name, f, "audio/ogg")},
            data={"model": "whisper-large-v3", "language": "es"},
            timeout=30
        )
    data = response.json()
    return data.get("text", "").strip()

# ── ElevenLabs texto a voz ────────────────────────────────────────────────────
def text_to_speech(text):
    try:
        from gtts import gTTS
        tts = gTTS(text=text, lang='es', slow=False)
        path = f"/tmp/jarvis_voice_{datetime.now().strftime('%H%M%S')}.mp3"
        tts.save(path)
        return path
    except Exception as e:
        logger.error(f"gTTS error: {e}")
        return None

# ── Imagen ────────────────────────────────────────────────────────────────────
def generate_image(prompt):
    try:
        HF_TOKEN = os.environ.get("HF_TOKEN", "")
        API_URL = "https://api-inference.huggingface.co/models/black-forest-labs/FLUX.1-schnell"
        headers = {"Authorization": f"Bearer {HF_TOKEN}"}
        response = requests.post(
            API_URL,
            headers=headers,
            json={"inputs": prompt},
            timeout=60
        )
        logger.info(f"HF status: {response.status_code}")
        if response.status_code == 200:
            path = f"/tmp/jarvis_img_{datetime.now().strftime('%H%M%S')}.jpg"
            Path(path).write_bytes(response.content)
            return path
        logger.error(f"HF error: {response.text[:200]}")
    except Exception as e:
        logger.error(f"Image error: {e}")
    return None

# ── Crear Word ────────────────────────────────────────────────────────────────
def create_word_doc(content, filename="documento.docx"):
    try:
        from docx import Document
        doc = Document()
        for line in content.split("\n"):
            line = line.strip()
            if not line:
                doc.add_paragraph()
            elif line.startswith("## "):
                doc.add_heading(line[3:], level=2)
            elif line.startswith("# "):
                doc.add_heading(line[2:], level=1)
            elif line.startswith("- "):
                doc.add_paragraph(line[2:], style="List Bullet")
            else:
                doc.add_paragraph(line)
        path = f"/tmp/{filename}"
        doc.save(path)
        return path
    except Exception as e:
        logger.error(f"Word error: {e}")
        return None

# ── Crear Excel ───────────────────────────────────────────────────────────────
def create_excel(csv_content, filename="datos.xlsx"):
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Datos"
        lines = [l for l in csv_content.strip().split("\n") if l.strip()]
        for row_idx, line in enumerate(lines, 1):
            cells = [c.strip().strip('"') for c in line.split(",")]
            for col_idx, val in enumerate(cells, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=val)
                if row_idx == 1:
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill("solid", fgColor="1D9E75")
                    cell.alignment = Alignment(horizontal="center")
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=10)
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 40)
        path = f"/tmp/{filename}"
        wb.save(path)
        return path
    except Exception as e:
        logger.error(f"Excel error: {e}")
        return None

# ── Crear PowerPoint ──────────────────────────────────────────────────────────
def create_pptx(slides_json, filename="presentacion.pptx"):
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        slides_data = json.loads(slides_json)
        prs = Presentation()
        for slide_data in slides_data:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = slide_data.get("titulo", "")
            title.text_frame.paragraphs[0].font.size = Pt(32)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1D, 0x9E, 0x75)
            tf = body.text_frame
            tf.clear()
            for i, punto in enumerate(slide_data.get("puntos", [])):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = punto
                p.font.size = Pt(20)
        path = f"/tmp/{filename}"
        prs.save(path)
        return path
    except Exception as e:
        logger.error(f"PPTX error: {e}")
        return None

# ── Procesar mensaje de texto o voz ──────────────────────────────────────────
def process_message(user_id, text, respond_with_voice=False):
    agent = detect_agent(text)
    user_mem = get_user_memory(user_id)

    try:
        reply = ask_groq(agent, text, user_mem.get("history", []), user_mem.get("items", []))
    except Exception as e:
        logger.error(f"Groq error: {e}")
        return None, None, "error"

    user_mem.setdefault("history", [])
    user_mem["history"].append({"role": "user", "content": text})
    user_mem["history"].append({"role": "assistant", "content": reply})
    user_mem["history"] = user_mem["history"][-40:]

    clean_reply, mem_item = extract_memory_tag(reply)
    if mem_item:
        user_mem.setdefault("items", []).append(mem_item)
    update_user_memory(user_id, user_mem)

    return clean_reply, agent, reply

# ── Handlers ──────────────────────────────────────────────────────────────────
@bot.message_handler(commands=["start"])
def start(message):
    bot.reply_to(message,
        "👋 Hola Eduardo, soy Jarvis.\n\n"
        "Puedo ayudarte con:\n"
        "🖼 Imágenes — 'genera una imagen de...'\n"
        "📄 Word — 'crea un documento sobre...'\n"
        "📊 Excel — 'haz una tabla con...'\n"
        "📊 PowerPoint — 'crea una presentación de...'\n"
        "🎤 Mensajes de voz — mándame un audio\n"
        "💬 Cualquier pregunta\n\n"
        "¿En qué te ayudo?"
    )

@bot.message_handler(content_types=["voice"])
def handle_voice(message):
    bot.send_chat_action(message.chat.id, "typing")
    user_id = message.from_user.id

    # Descargar audio
    file_info = bot.get_file(message.voice.file_id)
    downloaded = bot.download_file(file_info.file_path)
    with tempfile.NamedTemporaryFile(delete=False, suffix=".ogg") as f:
        f.write(downloaded)
        tmp_path = f.name

    # Transcribir
    try:
        text = transcribe_audio(tmp_path)
        Path(tmp_path).unlink(missing_ok=True)
    except Exception as e:
        logger.error(f"Transcription error: {e}")
        bot.reply_to(message, "No pude entender el audio. Intenta de nuevo.")
        return

    if not text:
        bot.reply_to(message, "No entendí el audio. Intenta hablar más claro.")
        return

    bot.reply_to(message, f"🎤 Escuché: _{text}_", parse_mode="Markdown")
    bot.send_chat_action(message.chat.id, "typing")

    clean_reply, agent, raw_reply = process_message(user_id, text, respond_with_voice=True)
    if not clean_reply:
        bot.reply_to(message, "Error procesando tu mensaje.")
        return

    if agent == "imagen":
        bot.send_chat_action(message.chat.id, "upload_photo")
        path = generate_image(raw_reply)
        if path:
            with open(path, "rb") as f:
                bot.send_photo(message.chat.id, f, caption=f"Prompt: {raw_reply[:200]}")
            Path(path).unlink(missing_ok=True)
        else:
            bot.reply_to(message, "No pude generar la imagen.")
        return

    # Responder con voz
    bot.send_chat_action(message.chat.id, "record_audio")
    voice_path = text_to_speech(clean_reply)
    if voice_path:
        with open(voice_path, "rb") as f:
            bot.send_voice(message.chat.id, f)
        Path(voice_path).unlink(missing_ok=True)
    else:
        bot.reply_to(message, clean_reply)

@bot.message_handler(content_types=["text"])
def handle_message(message):
    user_id = message.from_user.id
    text = message.text or ""
    agent = detect_agent(text)
    user_mem = get_user_memory(user_id)

    bot.send_chat_action(message.chat.id, "typing")

    try:
        reply = ask_groq(agent, text, user_mem.get("history", []), user_mem.get("items", []))
    except Exception as e:
        logger.error(f"Groq error: {e}")
        bot.reply_to(message, "Error al contactar la IA. Intenta de nuevo.")
        return

    user_mem.setdefault("history", [])
    user_mem["history"].append({"role": "user", "content": text})
    user_mem["history"].append({"role": "assistant", "content": reply})
    user_mem["history"] = user_mem["history"][-40:]

    if agent == "imagen":
        bot.send_chat_action(message.chat.id, "upload_photo")
        path = generate_image(reply)
        if path:
            with open(path, "rb") as f:
                bot.send_photo(message.chat.id, f, caption=f"🎨 {reply[:200]}")
            Path(path).unlink(missing_ok=True)
        else:
            bot.reply_to(message, "No pude generar la imagen. Intenta con otra descripción.")
        update_user_memory(user_id, user_mem)
        return

    if agent == "word":
        clean_reply, mem_item = extract_memory_tag(reply)
        fname = f"Jarvis_{datetime.now().strftime('%Y%m%d_%H%M')}.docx"
        path = create_word_doc(clean_reply, fname)
        if path:
            with open(path, "rb") as f:
                bot.send_document(message.chat.id, f, caption="📄 Aquí está tu documento Word.")
            Path(path).unlink(missing_ok=True)
        else:
            bot.reply_to(message, clean_reply)
        if mem_item:
            user_mem.setdefault("items", []).append(mem_item)
        update_user_memory(user_id, user_mem)
        return

    if agent == "excel":
        fname = f"Jarvis_{datetime.now().strftime('%Y%m%d_%H%M')}.xlsx"
        path = create_excel(reply, fname)
        if path:
            with open(path, "rb") as f:
                bot.send_document(message.chat.id, f, caption="📊 Aquí está tu Excel.")
            Path(path).unlink(missing_ok=True)
        else:
            bot.reply_to(message, reply)
        update_user_memory(user_id, user_mem)
        return

    if agent == "pptx":
        fname = f"Jarvis_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
        path = create_pptx(reply, fname)
        if path:
            with open(path, "rb") as f:
                bot.send_document(message.chat.id, f, caption="📊 Aquí está tu presentación.")
            Path(path).unlink(missing_ok=True)
        else:
            bot.reply_to(message, reply)
        update_user_memory(user_id, user_mem)
        return

    clean_reply, mem_item = extract_memory_tag(reply)
    if mem_item:
        user_mem.setdefault("items", []).append(mem_item)
    update_user_memory(user_id, user_mem)
    bot.reply_to(message, clean_reply)

@bot.message_handler(content_types=["document"])
def handle_document(message):
    doc = message.document
    file_info = bot.get_file(doc.file_id)
    downloaded = bot.download_file(file_info.file_path)
    suffix = Path(doc.file_name).suffix
    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as f:
        f.write(downloaded)
        tmp_path = f.name

    bot.send_chat_action(message.chat.id, "typing")
    content = ""
    try:
        if doc.file_name.endswith(".txt"):
            content = Path(tmp_path).read_text(errors="ignore")[:4000]
        elif doc.file_name.endswith(".pdf"):
            import pdfplumber
            with pdfplumber.open(tmp_path) as pdf:
                content = "\n".join(p.extract_text() or "" for p in pdf.pages)[:4000]
        elif doc.file_name.endswith(".docx"):
            from docx import Document
            d = Document(tmp_path)
            content = "\n".join(p.text for p in d.paragraphs)[:4000]
    except Exception as e:
        logger.error(f"Doc error: {e}")

    Path(tmp_path).unlink(missing_ok=True)

    if not content:
        bot.reply_to(message, "Recibí el archivo pero no pude leerlo. Intenta con .txt, .pdf o .docx.")
        return

    user_id = message.from_user.id
    user_mem = get_user_memory(user_id)
    prompt = f"El usuario envió '{doc.file_name}'. Contenido:\n\n{content}\n\nResúmelo y pregunta qué quiere hacer."
    try:
        reply = ask_groq("director", prompt, user_mem.get("history", []), user_mem.get("items", []))
        clean_reply, _ = extract_memory_tag(reply)
        bot.reply_to(message, clean_reply)
    except Exception as e:
        bot.reply_to(message, "Error procesando el documento.")

# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    logger.info("Jarvis iniciado ✅")
    bot.infinity_polling()
